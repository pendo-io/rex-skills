#!/usr/bin/env python3
"""Build a .pptx from a JSON spec, or extract a .pptx back to JSON.

Spec schema documented in ../SKILL.md. Original code, not derivative of any
proprietary skill — uses python-pptx (BSD-3) directly.
"""

from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Emu, Inches, Pt
except ImportError:
    print("Missing dep: pip install python-pptx", file=sys.stderr)
    sys.exit(2)


THEMES = {
    "pendo": {
        "accent": RGBColor(0xEB, 0x00, 0x8B),
        "ink": RGBColor(0x11, 0x11, 0x11),
        "muted": RGBColor(0x6B, 0x6B, 0x6B),
        "bg": RGBColor(0xFF, 0xFF, 0xFF),
        "title_font": "Calibri",
        "body_font": "Calibri",
    },
    "mono": {
        "accent": RGBColor(0x00, 0x00, 0x00),
        "ink": RGBColor(0x11, 0x11, 0x11),
        "muted": RGBColor(0x55, 0x55, 0x55),
        "bg": RGBColor(0xFF, 0xFF, 0xFF),
        "title_font": "Helvetica",
        "body_font": "Helvetica",
    },
}

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _set_slide_size(prs: Presentation) -> None:
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H


def _blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank


def _add_textbox(slide, left, top, width, height, text, *, font, size, bold=False, color=None, align="left"):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if align == "center":
        from pptx.enum.text import PP_ALIGN
        p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color
    return box


def _add_accent_bar(slide, theme):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(1.4), Inches(0.6), Emu(38100))  # ~3pt
    bar.fill.solid()
    bar.fill.fore_color.rgb = theme["accent"]
    bar.line.fill.background()


def _add_footer(slide, theme, text):
    if not text:
        return
    _add_textbox(
        slide, Inches(0.75), Inches(7.05), Inches(11.8), Inches(0.3),
        text, font=theme["body_font"], size=10, color=theme["muted"],
    )


def _title_slide(prs, theme, slide):
    s = _blank_slide(prs)
    _add_textbox(
        s, Inches(0.75), Inches(2.6), Inches(12), Inches(1.4),
        slide.get("title", ""), font=theme["title_font"], size=44, bold=True,
        color=theme["ink"],
    )
    if slide.get("subtitle"):
        _add_textbox(
            s, Inches(0.75), Inches(4.0), Inches(12), Inches(0.8),
            slide["subtitle"], font=theme["body_font"], size=20, color=theme["muted"],
        )
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(2.4), Inches(1.5), Emu(57150))
    bar.fill.solid(); bar.fill.fore_color.rgb = theme["accent"]; bar.line.fill.background()


def _section_slide(prs, theme, slide):
    s = _blank_slide(prs)
    _add_textbox(
        s, Inches(0.75), Inches(3.0), Inches(12), Inches(1.5),
        slide.get("title", ""), font=theme["title_font"], size=40, bold=True,
        color=theme["accent"],
    )


def _heading(slide_obj, theme, title):
    _add_textbox(
        slide_obj, Inches(0.75), Inches(0.6), Inches(12), Inches(0.7),
        title, font=theme["title_font"], size=26, bold=True, color=theme["ink"],
    )
    bar = slide_obj.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(1.35), Inches(0.6), Emu(19050),
    )
    bar.fill.solid(); bar.fill.fore_color.rgb = theme["accent"]; bar.line.fill.background()


def _bullets_slide(prs, theme, slide):
    s = _blank_slide(prs)
    _heading(s, theme, slide.get("title", ""))
    bullets = slide.get("bullets") or []
    box = s.shapes.add_textbox(Inches(0.75), Inches(1.7), Inches(11.8), Inches(5.0))
    tf = box.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
        p.level = 0
        for run in p.runs:
            run.font.name = theme["body_font"]
            run.font.size = Pt(20)
            run.font.color.rgb = theme["ink"]


def _two_column_slide(prs, theme, slide):
    s = _blank_slide(prs)
    _heading(s, theme, slide.get("title", ""))
    for col_key, x in (("left", Inches(0.75)), ("right", Inches(7.0))):
        col = slide.get(col_key) or {}
        _add_textbox(
            s, x, Inches(1.7), Inches(5.8), Inches(0.5),
            col.get("heading", ""), font=theme["title_font"], size=16, bold=True,
            color=theme["accent"],
        )
        box = s.shapes.add_textbox(x, Inches(2.2), Inches(5.8), Inches(4.6))
        tf = box.text_frame; tf.word_wrap = True
        for i, b in enumerate(col.get("bullets") or []):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = b
            for run in p.runs:
                run.font.name = theme["body_font"]
                run.font.size = Pt(16)
                run.font.color.rgb = theme["ink"]


def _quote_slide(prs, theme, slide):
    s = _blank_slide(prs)
    _add_textbox(
        s, Inches(1.5), Inches(2.5), Inches(10.3), Inches(2.5),
        f"“{slide.get('quote', '')}”",
        font=theme["title_font"], size=32, bold=False, color=theme["ink"],
    )
    if slide.get("attribution"):
        _add_textbox(
            s, Inches(1.5), Inches(5.2), Inches(10.3), Inches(0.6),
            f"— {slide['attribution']}", font=theme["body_font"], size=16,
            color=theme["muted"],
        )


def _image_slide(prs, theme, slide):
    s = _blank_slide(prs)
    _heading(s, theme, slide.get("title", ""))
    img = slide.get("image")
    if img and Path(img).exists():
        s.shapes.add_picture(img, Inches(2.0), Inches(1.8), height=Inches(5.2))
    else:
        _add_textbox(
            s, Inches(0.75), Inches(3.5), Inches(12), Inches(0.5),
            f"[image not found: {img!r}]", font=theme["body_font"], size=14,
            color=theme["muted"],
        )


def _table_slide(prs, theme, slide):
    s = _blank_slide(prs)
    _heading(s, theme, slide.get("title", ""))
    headers = slide.get("headers") or []
    rows = slide.get("rows") or []
    if not headers or not rows:
        return
    cols = len(headers)
    n_rows = len(rows) + 1
    tbl = s.shapes.add_table(
        n_rows, cols,
        Inches(0.75), Inches(1.8), Inches(11.8), Inches(0.5 * n_rows),
    ).table
    for c, h in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.text = str(h)
        for p in cell.text_frame.paragraphs:
            for run in p.runs:
                run.font.bold = True
                run.font.name = theme["title_font"]
                run.font.size = Pt(14)
                run.font.color.rgb = theme["bg"]
        cell.fill.solid(); cell.fill.fore_color.rgb = theme["accent"]
    for r, row in enumerate(rows, start=1):
        for c, val in enumerate(row[:cols]):
            cell = tbl.cell(r, c)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.name = theme["body_font"]
                    run.font.size = Pt(13)
                    run.font.color.rgb = theme["ink"]


SLIDE_BUILDERS = {
    "title": _title_slide,
    "section": _section_slide,
    "bullets": _bullets_slide,
    "two-column": _two_column_slide,
    "quote": _quote_slide,
    "image": _image_slide,
    "table": _table_slide,
}


def build_pptx(spec: dict, output: Path) -> Path:
    theme_name = spec.get("theme", "pendo")
    theme = THEMES.get(theme_name, THEMES["pendo"])
    prs = Presentation()
    _set_slide_size(prs)

    footer = " · ".join(filter(None, [spec.get("title"), spec.get("author")]))
    for slide_spec in spec.get("slides") or []:
        kind = slide_spec.get("type", "bullets")
        builder = SLIDE_BUILDERS.get(kind)
        if not builder:
            print(f"warn: unknown slide type: {kind}", file=sys.stderr)
            continue
        builder(prs, theme, slide_spec)
        if footer and kind != "title":
            _add_footer(prs.slides[-1], theme, footer)

    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output))
    return output


def extract_pptx(path: Path) -> dict:
    """Reverse the round-trip: read a .pptx into a slim JSON spec.

    Lossy by design — preserves slide titles, bullets, and quote-like text.
    """
    prs = Presentation(str(path))
    slides = []
    for s in prs.slides:
        title = ""
        bullets = []
        for shape in s.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if not text:
                continue
            if not title:
                title = text.split("\n")[0]
                rest = text.split("\n")[1:]
                bullets.extend([line for line in rest if line.strip()])
            else:
                bullets.extend([line for line in text.split("\n") if line.strip()])
        kind = "title" if not bullets else "bullets"
        entry: dict = {"type": kind, "title": title}
        if bullets:
            entry["bullets"] = bullets
        slides.append(entry)
    return {"title": Path(path).stem, "theme": "pendo", "slides": slides}


def main() -> int:
    ap = argparse.ArgumentParser(description="Deck Builder: JSON ↔ .pptx")
    ap.add_argument("--spec", help="Input JSON spec file")
    ap.add_argument("--output", help="Output .pptx (build mode) or .json (extract mode)")
    ap.add_argument("--extract", help="Existing .pptx to convert back to JSON")
    args = ap.parse_args()

    if args.extract:
        spec = extract_pptx(Path(args.extract))
        out = Path(args.output or "deck.json")
        out.write_text(json.dumps(spec, indent=2))
        print(f"wrote {out}")
        return 0

    if not args.spec:
        ap.error("either --spec or --extract is required")

    spec_text = Path(args.spec).read_text() if Path(args.spec).exists() else args.spec
    spec = json.loads(spec_text)
    out = Path(args.output or "deck.pptx")
    build_pptx(spec, out)
    print(f"wrote {out}  ({len(spec.get('slides') or [])} slides)")
    return 0


if __name__ == "__main__":
    sys.exit(main())
